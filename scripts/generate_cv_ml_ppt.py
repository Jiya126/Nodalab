#!/usr/bin/env python3
"""Generate Computer Vision & ML for Robotics lecture deck (reference-style)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (matches Sensors deck tone) ─────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x7A, 0x8C)
SLATE = RGBColor(0x4A, 0x55, 0x68)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE8, 0x6C, 0x00)
MUTED = RGBColor(0x71, 0x85, 0x94)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.75)
MARGIN_R = Inches(0.75)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


def _set_run(run, size, bold=False, color=NAVY, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _footer_bar(slide):
    _add_rect(slide, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), NAVY)


def _slide_number(slide, num):
    box = slide.shapes.add_textbox(SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.55), Inches(0.6), Inches(0.3))
    tf = box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run = tf.paragraphs[0].add_run()
    run.text = str(num)
    _set_run(run, 11, color=WHITE)


def _tagline(slide, text, slide_num):
    box = slide.shapes.add_textbox(MARGIN_L, SLIDE_H - Inches(1.05), CONTENT_W, Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, 13, bold=True, color=TEAL)
    _footer_bar(slide)
    _slide_number(slide, slide_num)


def _category_label(slide, label):
    box = slide.shapes.add_textbox(MARGIN_L, Inches(0.55), CONTENT_W, Inches(0.35))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = label.upper()
    _set_run(run, 11, bold=True, color=TEAL)


def _title(slide, text, top=Inches(0.95)):
    box = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    _set_run(run, 28, bold=True, color=NAVY)


def _bullets(slide, items, top=Inches(1.85), width=None, size=18):
    w = width or CONTENT_W
    box = slide.shapes.add_textbox(MARGIN_L, top, w, SLIDE_H - top - Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item[0], item[1]
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        _set_run(run, size if level == 0 else size - 2, color=SLATE if level == 0 else MUTED)


def _content_slide(prs, label, title, bullets, tagline, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), TEAL)
    _category_label(slide, label)
    _title(slide, title)
    _bullets(slide, bullets)
    _tagline(slide, tagline, num)
    return slide


def _section_divider(prs, num_str, title, subtitle, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    _add_rect(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06), TEAL)

    num_box = slide.shapes.add_textbox(MARGIN_L, Inches(2.0), CONTENT_W, Inches(0.8))
    run = num_box.text_frame.paragraphs[0].add_run()
    run.text = num_str
    _set_run(run, 48, bold=True, color=TEAL)

    tbox = slide.shapes.add_textbox(MARGIN_L, Inches(3.45), CONTENT_W, Inches(1.2))
    run = tbox.text_frame.paragraphs[0].add_run()
    run.text = title
    _set_run(run, 36, bold=True, color=WHITE)

    sbox = slide.shapes.add_textbox(MARGIN_L, Inches(4.55), CONTENT_W, Inches(0.6))
    run = sbox.text_frame.paragraphs[0].add_run()
    run.text = subtitle
    _set_run(run, 18, color=RGBColor(0xA0, 0xAE, 0xC0))

    _slide_number(slide, slide_num)
    return slide


def _title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    _add_rect(slide, Inches(0), Inches(5.85), SLIDE_W, Inches(0.08), TEAL)

    series = slide.shapes.add_textbox(MARGIN_L, Inches(1.2), CONTENT_W, Inches(0.5))
    run = series.text_frame.paragraphs[0].add_run()
    run.text = "COMPUTER VISION & MACHINE LEARNING FOR ROBOTICS"
    _set_run(run, 14, bold=True, color=TEAL)

    main = slide.shapes.add_textbox(MARGIN_L, Inches(2.0), CONTENT_W, Inches(1.4))
    tf = main.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = "Computer Vision & ML in Automation"
    _set_run(run, 40, bold=True, color=WHITE)

    topics = slide.shapes.add_textbox(MARGIN_L, Inches(3.55), CONTENT_W, Inches(0.5))
    run = topics.text_frame.paragraphs[0].add_run()
    run.text = "OpenCV · Edge & Object Detection · YOLO · Stereo Vision · Depth · ML Deployment"
    _set_run(run, 16, color=RGBColor(0xA0, 0xAE, 0xC0))

    sub = slide.shapes.add_textbox(MARGIN_L, Inches(4.15), CONTENT_W, Inches(0.5))
    run = sub.text_frame.paragraphs[0].add_run()
    run.text = "Fundamentals · Practical ML · Real-Time Detection Lab"
    _set_run(run, 16, color=RGBColor(0xA0, 0xAE, 0xC0))

    deck = slide.shapes.add_textbox(MARGIN_L, Inches(6.15), CONTENT_W, Inches(0.4))
    run = deck.text_frame.paragraphs[0].add_run()
    run.text = "Mixed Module · 40-Slide Lecture Deck"
    _set_run(run, 13, bold=True, color=ACCENT)

    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    n = 1
    _title_slide(prs)
    n += 1

    # ── Roadmap ───────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), TEAL)
    _category_label(slide, "Course Roadmap")
    _title(slide, "Agenda")
    _bullets(
        slide,
        [
            "Computer vision fundamentals",
            ("OpenCV pipeline & image representation", 1),
            ("Edge detection, features & classical detection", 1),
            ("YOLO basics, stereo vision & depth estimation", 1),
            "ML in automation — practical overview",
            ("Supervised learning for fault detection", 1),
            ("Reinforcement learning for path planning", 1),
            ("Model deployment on edge & factory floor", 1),
            "Lab: object detection mini project",
            ("Webcam pipeline · color-based & YOLO-lite detection", 1),
        ],
        top=Inches(1.75),
        size=17,
    )
    _tagline(slide, "From pixels to decisions — perception, learning, and deployment in robotics.", n)
    n += 1

    slides_data = [
        (
            "INTRODUCTION",
            "Why Computer Vision in Robotics?",
            [
                "Robots must perceive unstructured environments — vision is the richest sensor modality.",
                ("Cameras are passive, low-cost, and information-dense compared to LiDAR or ultrasonic.", 1),
                ("CV converts images into actionable data: objects, poses, obstacles, defects.", 1),
                "Modern pipelines combine classical geometry with deep learning.",
                ("Classical — edges, features, stereo, calibration.", 1),
                ("Deep learning — detection, segmentation, tracking at scale.", 1),
                "Latency, lighting, and compute budget shape every design choice.",
                "This module bridges theory (CV + ML) with a hands-on detection lab.",
            ],
            "Vision turns raw pixels into the semantic world model a robot can act on.",
        ),
        (
            "SYSTEMS VIEW",
            "Vision in the Perception Stack",
            [
                "Typical robot perception pipeline:",
                ("Acquire → preprocess → detect/segment → track → fuse → plan → act", 1),
                "Camera feeds enter alongside encoders, IMU, and range sensors.",
                ("Multi-sensor fusion improves robustness when vision alone fails.", 1),
                "Frame rate & latency must match control loop requirements.",
                ("30 fps capture ≠ 30 Hz closed-loop if inference takes 80 ms.", 1),
                "Edge deployment avoids cloud round-trip but limits model size.",
                "Safety-critical systems need deterministic fallbacks when models fail.",
            ],
            "Perception is part of the control loop — latency and reliability matter as much as accuracy.",
        ),
        (
            "FOUNDATIONS",
            "Digital Images & OpenCV Basics",
            [
                "An image is a 2D array of pixels — each pixel holds intensity or colour channels.",
                ("Grayscale: single channel (0–255). Colour: typically BGR in OpenCV.", 1),
                "OpenCV (Open Source Computer Vision Library) is the de-facto CV toolkit.",
                ("Read/write: cv2.imread(), cv2.imwrite()", 1),
                ("Display: cv2.imshow() · Convert colour: cv2.cvtColor()", 1),
                ("Resize, crop, rotate, threshold — core pre-processing ops", 1),
                "Pre-processing reduces noise and normalises input before higher-level tasks.",
                "Always measure resolution, fps, and memory when targeting embedded hardware.",
            ],
            "OpenCV is the workbench — master I/O and pre-processing before detection.",
        ),
        (
            "FOUNDATIONS",
            "Image Pre-Processing & Filtering",
            [
                "Raw images contain sensor noise, uneven lighting, and motion blur.",
                "Gaussian blur — smooths high-frequency noise before edge extraction.",
                ("Kernel size controls smoothing strength vs detail loss.", 1),
                "Histogram equalisation — boosts contrast in low-light scenes.",
                "Morphological ops — erode/dilate/open/close for shape: blobs & holes.",
                "Colour spaces beyond BGR:",
                ("HSV — separates hue from lighting; ideal for colour tracking.", 1),
                ("Lab — perceptually uniform; useful for consistent thresholds.", 1),
                "Correct pre-processing often matters more than a fancier detector.",
            ],
            "Clean, normalised images make every downstream algorithm more reliable.",
        ),
    ]

    for label, title, bullets, tag in slides_data:
        _content_slide(prs, label, title, bullets, tag, n)
        n += 1

    _section_divider(prs, "06", "Computer Vision Fundamentals", "OpenCV · Edges · Detection · YOLO · Stereo · Depth", n)
    n += 1

    cv_slides = [
        (
            "EDGE DETECTION",
            "Finding Structure: Edges & Gradients",
            [
                "Edges mark boundaries where intensity changes sharply — corners, object outlines.",
                "Gradient magnitude & direction computed via Sobel or Scharr operators.",
                ("Sobel approximates first derivative in x and y.", 1),
                "Canny edge detector — industry standard multi-stage pipeline:",
                ("1. Gaussian blur  2. Gradient  3. Non-max suppression  4. Hysteresis threshold", 1),
                "Parameters: low/high thresholds control edge fragmentation vs noise.",
                "Applications: lane detection, part localisation, visual servoing features.",
                "Limitation: edges alone do not identify objects — they provide geometry hints.",
            ],
            "Canny gives thin, well-localised edges — the starting point for many CV pipelines.",
        ),
        (
            "FEATURES",
            "Keypoints, Descriptors & Classical Detection",
            [
                "Keypoints — repeatable interest points (corners, blobs) across frames.",
                "Harris corner detector — responds to intensity changes in multiple directions.",
                "ORB / SIFT / SURF — detect + describe patches for matching & tracking.",
                ("ORB: fast, binary, patent-free — good for embedded robotics.", 1),
                "Haar cascades & HOG + SVM — classical object detection before deep learning.",
                ("Viola–Jones for faces; HOG for pedestrians.", 1),
                "Classical methods are lightweight but struggle with scale, pose, and clutter.",
                "Still valuable when compute is tight or training data is scarce.",
            ],
            "Hand-crafted features remain useful when deep models are too heavy or opaque.",
        ),
        (
            "OBJECT DETECTION",
            "From Classification to Localisation",
            [
                "Classification — what object?  Detection — what + where (bounding box)?",
                "Two-stage detectors (R-CNN family): region proposals → classify each region.",
                ("High accuracy, slower inference — less common on mobile robots today.", 1),
                "One-stage detectors: predict boxes + classes in a single forward pass.",
                ("YOLO, SSD, RetinaNet — optimised for real-time performance.", 1),
                "Metrics that matter:",
                ("IoU (Intersection over Union) — box overlap quality.", 1),
                ("Precision / recall / mAP — trade-off between false alarms and misses.", 1),
                "Robotics needs consistent latency, not just peak accuracy on a benchmark.",
            ],
            "Detection = simultaneous recognition and localisation — the core of robot perception.",
        ),
        (
            "YOLO",
            "YOLO Basics: You Only Look Once",
            [
                "YOLO divides the image into an S×S grid; each cell predicts boxes + class scores.",
                "Single network pass → extremely fast compared to two-stage pipelines.",
                "Versions: YOLOv5/v8 (Ultralytics) popular for prototyping and edge export.",
                "Key concepts:",
                ("Anchor boxes — prior shapes the network refines.", 1),
                ("Non-max suppression (NMS) — removes duplicate overlapping detections.", 1),
                ("Confidence score — how sure the model is an object exists.", 1),
                "Training requires labelled datasets (COCO, custom factory images).",
                "YOLO-lite / nano variants trade accuracy for speed on Raspberry Pi & Jetson.",
            ],
            "YOLO trades a little accuracy for the real-time speed robots need.",
        ),
        (
            "STEREO VISION",
            "Binocular Geometry & Disparity",
            [
                "Two calibrated cameras separated by baseline B observe the same scene.",
                "Disparity d = x_left − x_right for corresponding points.",
                "Depth Z = (f × B) / d  — f = focal length in pixels.",
                "Stereo matching finds correspondences:",
                ("Block matching (BM) — fast, less accurate.", 1),
                ("Semi-global block matching (SGBM) — smoother disparity maps.", 1),
                "Challenges: textureless regions, occlusions, reflective surfaces.",
                "Calibration (intrinsics + extrinsics) is mandatory for metric depth.",
                "Used in pick-and-place, obstacle avoidance, and mobile manipulation.",
            ],
            "Stereo turns a pair of cheap cameras into a 3D range sensor.",
        ),
        (
            "DEPTH",
            "Depth Estimation Beyond Stereo",
            [
                "Structured light (Kinect-style) — projected pattern + triangulation.",
                "Time-of-Flight (ToF) — measures round-trip of modulated IR light.",
                "LiDAR — sparse/dense 3D point clouds; complements RGB cameras.",
                "Monocular depth estimation — deep networks infer depth from a single image.",
                ("Useful when only one camera is available; scale may be ambiguous.", 1),
                "Sensor fusion: RGB-D + IMU + wheel odometry for robust localisation.",
                "Trade-offs: range, resolution, outdoor performance, cost, power.",
                "Choose depth modality based on environment — factory floor vs outdoor AGV.",
            ],
            "No single depth sensor wins everywhere — match the modality to the task.",
        ),
        (
            "PIPELINE",
            "End-to-End Vision Processing Pipeline",
            [
                "Acquisition → undistort → resize/normalise → infer → post-process → output.",
                "Camera calibration removes lens distortion for accurate geometry.",
                ("Checkerboard calibration → camera matrix + distortion coefficients.", 1),
                "Coordinate transforms link camera frame → robot base → world frame.",
                "Tracking (Kalman, SORT, DeepSORT) maintains identity across frames.",
                "Visual servoing closes the loop: image error → joint/cartesian velocity.",
                "Log latency at each stage — bottlenecks are often pre-processing, not the model.",
            ],
            "A production pipeline is more than inference — geometry and timing complete the story.",
        ),
        (
            "TRACKING",
            "Optical Flow & Multi-Object Tracking",
            [
                "Optical flow estimates per-pixel motion between consecutive frames.",
                ("Lucas–Kanade — sparse flow on corner features.", 1),
                ("Dense Farneback — full motion field for scene analysis.", 1),
                "Tracking associates detections across frames to maintain object IDs.",
                ("SORT: Kalman predict + Hungarian match on IoU.", 1),
                ("DeepSORT adds appearance embedding for re-identification.", 1),
                "Critical for pick-and-place, conveyor counting, and follow-me robots.",
                "Trade-off: more robust tracking adds latency — profile on target hardware.",
            ],
            "Detection finds objects once; tracking keeps identity stable over time.",
        ),
        (
            "SEGMENTATION",
            "Semantic & Instance Segmentation",
            [
                "Detection draws boxes; segmentation labels every pixel.",
                "Semantic segmentation — each pixel gets a class (road, person, machine).",
                "Instance segmentation — separates individual object instances.",
                "Architectures: U-Net, DeepLab, Mask R-CNN, YOLO-seg variants.",
                "Robotics uses:",
                ("Grasp planning on segmented object masks.", 1),
                ("Free-space navigation from traversable-region maps.", 1),
                ("Quality inspection via pixel-level defect maps.", 1),
                "Heavier than bounding-box detection — use when shape matters, not just location.",
            ],
            "Segmentation gives pixel-accurate masks — essential for manipulation and inspection.",
        ),
    ]

    for label, title, bullets, tag in cv_slides:
        _content_slide(prs, label, title, bullets, tag, n)
        n += 1

    _section_divider(prs, "14", "ML in Automation", "Supervised Learning · Fault Detection · RL · Deployment", n)
    n += 1

    ml_slides = [
        (
            "ML OVERVIEW",
            "Machine Learning in Industrial Automation",
            [
                "ML learns patterns from data instead of relying on hand-written rules.",
                "Three paradigms most relevant to robotics & automation:",
                ("Supervised — labelled examples (classification, regression).", 1),
                ("Unsupervised — find structure without labels (clustering, anomaly).", 1),
                ("Reinforcement learning — agent learns via reward in an environment.", 1),
                "Factory constraints: labelled data cost, explainability, safety certification.",
                "Start with a clear metric — defect rate, cycle time, collision rate, uptime.",
                "Baseline with classical CV before jumping to deep learning.",
            ],
            "Pick the simplest learning paradigm that solves the problem with acceptable risk.",
        ),
        (
            "SUPERVISED",
            "Supervised Learning for Fault Detection",
            [
                "Goal: classify products or signals as OK vs defective (or multi-class faults).",
                "Workflow:",
                ("1. Collect & label images/traces  2. Split train/val/test", 1),
                ("3. Train classifier (CNN, transfer learning)  4. Evaluate on held-out set", 1),
                ("5. Deploy with confidence threshold & human review loop", 1),
                "Transfer learning — fine-tune ResNet/EfficientNet on small factory datasets.",
                "Data quality dominates: balanced classes, consistent lighting, label accuracy.",
                "Metrics: precision (avoid false rejects), recall (catch all defects).",
                "Integrate with PLC/MES — trigger reject actuator when confidence > threshold.",
            ],
            "Fault detection is the highest-ROI ML use case in many production lines.",
        ),
        (
            "SUPERVISED",
            "Dataset Design & Model Training Tips",
            [
                "Capture variability: lighting shifts, part orientation, camera vibration.",
                "Augmentation — rotation, flip, colour jitter, noise — improves generalisation.",
                "Avoid data leakage: test set must reflect future production conditions.",
                "Class imbalance — use oversampling, weighted loss, or focal loss.",
                "Monitor drift: retrain when new defect types or lighting changes appear.",
                "Explainability tools (Grad-CAM) help engineers trust and debug decisions.",
                "Version datasets and models — reproducibility is essential for audit trails.",
            ],
            "A small, well-curated dataset beats a large, noisy one in production.",
        ),
        (
            "REINFORCEMENT LEARNING",
            "RL for Path Planning & Navigation",
            [
                "Agent observes state s, takes action a, receives reward r, transitions to s'.",
                "Goal: learn policy π(a|s) that maximises cumulative reward over time.",
                "Path planning applications:",
                ("Mobile robots — avoid obstacles, minimise time/energy.", 1),
                ("Manipulators — learn joint trajectories in cluttered spaces.", 1),
                "Algorithms: Q-learning, DQN, PPO, SAC — model-free methods dominate sim-to-real.",
                "Simulation (Gazebo, Isaac Sim) enables millions of safe training episodes.",
                "Reward shaping is critical — sparse rewards make learning impractically slow.",
                "Sim-to-real gap: randomise textures, dynamics, sensor noise before deployment.",
            ],
            "RL excels when the optimal rule is unknown but success/failure is measurable.",
        ),
        (
            "REINFORCEMENT LEARNING",
            "RL Practical Considerations",
            [
                "Safety during exploration — constrain actions, use shielded policies.",
                "Hybrid approach: classical planner (A*, RRT) + RL for local refinement.",
                "Observation design: lidar scan, occupancy grid, or end-to-end camera input.",
                "Training time can be hours/days — budget compute accordingly.",
                "Evaluate in staged rollout: sim → fenced test cell → production.",
                "Fallback to rule-based behaviour if policy confidence drops.",
                "Not every navigation problem needs RL — geometric planners are often enough.",
            ],
            "Use RL when the environment is complex and hand-tuned rules break down.",
        ),
        (
            "DEPLOYMENT",
            "Model Deployment on Edge & Factory Floor",
            [
                "Training (GPU cloud/workstation) ≠ inference (edge PLC, Jetson, IPC).",
                "Optimisation pipeline:",
                ("Quantisation (INT8) · pruning · knowledge distillation · TensorRT / ONNX", 1),
                "Containerise with Docker for reproducible factory deployments.",
                "REST/gRPC or MQTT interfaces connect vision service to robot controller.",
                "Monitor in production: latency p95, throughput, false positive rate, GPU temp.",
                "CI/CD for models — automated retrain triggers on data drift detection.",
                "Regulatory & safety: document failure modes; maintain kill-switch & manual mode.",
            ],
            "A model is not deployed until it runs reliably at target latency on target hardware.",
        ),
        (
            "DEPLOYMENT",
            "MLOps & Integration Patterns",
            [
                "Edge inference node sits between camera and PLC/robot controller.",
                "Typical architecture: camera → inference server → JSON/Modbus → actuator.",
                "Shadow mode — run new model in parallel without affecting production.",
                "A/B testing — route fraction of lines to candidate model.",
                "Logging: store misclassified images for continuous improvement.",
                "Security: signed models, network segmentation, no open camera streams.",
                "Align with IT/OT teams on update windows and rollback procedures.",
            ],
            "Treat deployed models like production software — versioned, monitored, and reversible.",
        ),
        (
            "ANOMALY DETECTION",
            "Unsupervised & Semi-Supervised Fault Finding",
            [
                "When defect labels are scarce, learn what 'normal' looks like.",
                "Autoencoders — high reconstruction error flags anomalies.",
                "One-class SVM / isolation forest on hand-crafted or CNN features.",
                "Semi-supervised: train on OK samples only; flag outliers at inference.",
                "Useful for novel defect types not seen during labelling.",
                "Combine with supervised model: unsupervised pre-filter → CNN confirm.",
                "Monitor reconstruction-error distributions for drift over time.",
            ],
            "Anomaly detection fills the gap when you have plenty of good parts but few bad examples.",
        ),
    ]

    for label, title, bullets, tag in ml_slides:
        _content_slide(prs, label, title, bullets, tag, n)
        n += 1

    _section_divider(prs, "22", "Lab: Object Detection Mini Project", "Webcam · OpenCV · Python · Color & YOLO-lite", n)
    n += 1

    lab_slides = [
        (
            "LAB OVERVIEW",
            "Object Detection Mini Project — Goals",
            [
                "Build a real-time object detection demo using a webcam and Python.",
                "Two detection tracks (choose one or compare both):",
                ("Track A — colour-based detection in HSV space (no ML required).", 1),
                ("Track B — YOLO-lite / pre-trained nano model for class detection.", 1),
                "Learning outcomes:",
                ("Capture & display live video with OpenCV.", 1),
                ("Draw bounding boxes, labels, and FPS overlay.", 1),
                ("Understand latency trade-offs between classical and deep methods.", 1),
                "Deliverable: working script + short report on accuracy and frame rate.",
            ],
            "Hands-on practice cements the theory — you will see latency and accuracy live.",
        ),
        (
            "LAB SETUP",
            "Environment & Dependencies",
            [
                "Hardware: USB webcam, laptop or Jetson/RPi (optional).",
                "Software stack:",
                ("Python 3.10+ · opencv-python · numpy", 1),
                ("Track B add-on: ultralytics (YOLOv8n) or opencv-dnn + ONNX model", 1),
                "Setup commands:",
                ("pip install opencv-python numpy ultralytics", 1),
                "Verify camera index: try cv2.VideoCapture(0) — change index if needed.",
                "Test at reduced resolution (640×480) first for smoother frame rates.",
                "Use a well-lit scene; avoid backlighting behind target objects.",
            ],
            "Get the camera pipeline working before adding any detection logic.",
        ),
        (
            "LAB — TRACK A",
            "Colour-Based Detection (HSV Thresholding)",
            [
                "Convert frame BGR → HSV: hsv = cv2.cvtColor(frame, cv2.COLOR_BGR2HSV)",
                "Define colour range: lower = np.array([H,S,V]), upper = np.array([...])",
                "Create mask: mask = cv2.inRange(hsv, lower, upper)",
                "Clean mask: cv2.erode/dilate or morphologyEx to remove speckle noise.",
                "Find contours: contours, _ = cv2.findContours(mask, ...)",
                "Filter by area; draw cv2.boundingRect() on large enough blobs.",
                "Tune H,S,V with trackbars — lighting changes require re-tuning.",
                "Pros: fast (~60 fps), no GPU. Cons: fails with similar colours & shadows.",
            ],
            "Colour tracking is the fastest path to a real-time bounding box on constrained targets.",
        ),
        (
            "LAB — TRACK A",
            "Improving Colour Detection Robustness",
            [
                "Use adaptive background subtraction if objects move against static scene.",
                "Combine two colour ranges for bi-coloured targets (mask1 | mask2).",
                "Apply Gaussian blur before thresholding to suppress sensor noise.",
                "Reject false positives with aspect-ratio and minimum-area filters.",
                "Overlay FPS: fps = 1 / (time.time() - t_prev) on each frame.",
                "Record a short video clip to analyse failure cases offline.",
                "Document HSV values used — reproducibility matters for the lab report.",
            ],
            "Robust classical detection is 80% good thresholds and contour filtering.",
        ),
        (
            "LAB — TRACK B",
            "YOLO-lite Detection with Ultralytics",
            [
                "Load pre-trained nano model: model = YOLO('yolov8n.pt')",
                "Run inference on frame: results = model(frame, verbose=False)",
                "Visualise: annotated = results[0].plot()",
                "Filter classes: results = model(frame, classes=[0])  # person only, etc.",
                "Adjust confidence: model(frame, conf=0.5) to reduce false positives.",
                "Expect 15–40 fps on laptop CPU; 60+ fps on GPU/Jetson.",
                "Export to ONNX/TensorRT for embedded targets if needed.",
            ],
            "Pre-trained YOLO gives instant multi-class detection — ideal for general objects.",
        ),
        (
            "LAB — TRACK B",
            "Custom Labels & Fine-Tuning (Optional Extension)",
            [
                "Collect 50–200 images of your target object in the workspace.",
                "Label with Roboflow, CVAT, or LabelImg (YOLO format txt annotations).",
                "Fine-tune: model.train(data='data.yaml', epochs=50, imgsz=640)",
                "Evaluate mAP on validation split before live deployment.",
                "Compare fine-tuned vs off-the-shelf model on your specific objects.",
                "Watch for overfitting if dataset is tiny — lean on augmentation.",
                "Extension: deploy exported ONNX via cv2.dnn for OpenCV-only runtime.",
            ],
            "Fine-tuning closes the gap when generic COCO classes do not match your parts.",
        ),
        (
            "LAB PIPELINE",
            "Real-Time Loop Architecture",
            [
                "Standard capture loop:",
                ("cap = cv2.VideoCapture(0)", 1),
                ("while True: ret, frame = cap.read() → detect → annotate → imshow", 1),
                ("if cv2.waitKey(1) == ord('q'): break", 1),
                "Measure end-to-end latency — capture + inference + render.",
                "Optional: run inference every Nth frame; track between frames.",
                "Threading: capture thread + inference thread for smoother display.",
                "Release resources: cap.release(); cv2.destroyAllWindows()",
            ],
            "Structure the loop cleanly — profiling each stage reveals your real bottleneck.",
        ),
        (
            "LAB REPORT",
            "Evaluation & Deliverables",
            [
                "Report should include:",
                ("Method chosen (colour vs YOLO) and rationale.", 1),
                ("Average FPS and hardware used.", 1),
                ("Screenshots of successful detections and failure cases.", 1),
                ("Parameter table (HSV bounds or model conf/NMS settings).", 1),
                "Compare tracks: speed vs robustness vs setup effort.",
                "Discuss one scenario where your method fails and how to fix it.",
                "Optional demo video (30 s) showing live detection.",
            ],
            "Critical analysis of failures demonstrates deeper understanding than a perfect demo.",
        ),
        (
            "LAB TROUBLESHOOTING",
            "Common Issues & Fixes",
            [
                "Black screen — wrong camera index; try 0, 1, 2.",
                "Low FPS — reduce resolution; skip frames; use nano model.",
                "No detections — lower conf threshold; check lighting.",
                "False positives — raise conf; tighten NMS; filter by box area.",
                "Colour drift — re-tune HSV when lighting changes.",
                "Import errors — verify venv and package versions.",
                "Permission denied on camera — add user to video group (Linux).",
            ],
            "Most lab failures are environment issues — not broken algorithms.",
        ),
    ]

    for label, title, bullets, tag in lab_slides:
        _content_slide(prs, label, title, bullets, tag, n)
        n += 1

    synthesis = [
        (
            "SYNTHESIS",
            "Classical CV vs Deep Learning — When to Use What",
            [
                "Colour / edge / stereo — fast, interpretable, low data requirement.",
                "Deep detection — robust to variation, needs data & compute.",
                "Hybrid pipelines are common: classical pre-filter → DL on ROI.",
                "Factory defect inspection often: transfer-learned CNN + rule-based reject.",
                "Navigation often: geometric planner + RL or cost-map refinement.",
                "Always benchmark on target hardware, not just development laptop.",
                "Maintain a fallback when the learned model is uncertain.",
            ],
            "The best system combines the right tool at each stage — not the fanciest model everywhere.",
        ),
        (
            "ROBOTICS",
            "Real-World Integration Checklist",
            [
                "☑ Camera mounted rigidly; vibration damped; lens clean.",
                "☑ Calibration documented and periodically re-verified.",
                "☑ Lighting controlled or compensated (diffuse, consistent).",
                "☑ Inference latency < control period budget.",
                "☑ Fail-safe behaviour defined for camera/model outage.",
                "☑ Dataset represents seasonal and shift variations.",
                "☑ Model version, threshold, and deploy date logged.",
            ],
            "Integration failures are usually optical and operational — not algorithmic.",
        ),
        (
            "FUSION",
            "Multi-Sensor Fusion for Robust Perception",
            [
                "Combine vision with encoders, IMU, LiDAR, and ultrasonic.",
                "Late fusion — independent estimates merged at decision level.",
                "Early fusion — raw or feature-level combination before inference.",
                "Kalman / particle filters propagate uncertainty across sensors.",
                "Example: camera detects object class; LiDAR supplies metric distance.",
                "Redundancy: if vision fails (glare), range sensors maintain safety stop.",
                "Design for graceful degradation — never single point of failure.",
            ],
            "Fusion turns brittle single-sensor perception into robust situational awareness.",
        ),
        (
            "RESOURCES",
            "Further Reading & Tools",
            [
                "OpenCV docs: docs.opencv.org — tutorials and API reference.",
                "Ultralytics YOLO: docs.ultralytics.com — train, export, deploy.",
                "Roboflow — dataset labelling, augmentation, and model hosting.",
                "Open3D / PCL — point-cloud processing for depth & LiDAR.",
                "ONNX Runtime / TensorRT — optimised edge inference.",
                "Papers: YOLO (Redmon), StereoBM/SGBM (OpenCV), Mask R-CNN (He et al.).",
                "Datasets: COCO, ImageNet, MVTec AD (anomaly), custom factory captures.",
            ],
            "Build on open tools — the ecosystem around OpenCV and YOLO is production-ready.",
        ),
    ]

    for label, title, bullets, tag in synthesis:
        _content_slide(prs, label, title, bullets, tag, n)
        n += 1

    # Wrap-up
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), TEAL)
    _category_label(slide, "Wrap-Up")
    _title(slide, "Key Takeaways")
    _bullets(
        slide,
        [
            "OpenCV provides the foundation: I/O, filtering, edges, geometry, and display.",
            "Object detection localises and classifies — YOLO enables real-time performance.",
            "Stereo vision and depth sensors give robots metric 3D understanding.",
            "Supervised ML excels at fault detection when labelled data is available.",
            "RL can learn navigation policies but demands simulation, safety, and careful rewards.",
            "Deployment requires optimisation, monitoring, and rollback — not just training.",
            "The lab ties it together: webcam → detect → annotate → measure FPS.",
        ],
        top=Inches(1.75),
        size=17,
    )
    foot = slide.shapes.add_textbox(MARGIN_L, SLIDE_H - Inches(1.35), CONTENT_W, Inches(0.5))
    run = foot.text_frame.paragraphs[0].add_run()
    run.text = "Computer Vision & ML in Automation · Computer Vision & Machine Learning for Robotics"
    _set_run(run, 12, bold=True, color=TEAL)
    _footer_bar(slide)
    _slide_number(slide, n)

    out = "/home/jjzzzz/Downloads/Computer-Vision-ML-Robotics.pptx"
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    build()
