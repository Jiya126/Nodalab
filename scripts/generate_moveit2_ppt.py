#!/usr/bin/env python3
"""Generate a concise MoveIt2 basics educational deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x7A, 0x8C)
SLATE = RGBColor(0x4A, 0x55, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE8, 0x6C, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.9)
CONTENT_W = SLIDE_W - Inches(1.8)


def _run(run, size, bold=False, color=NAVY):
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()


def _footer(slide, tag, num):
    box = slide.shapes.add_textbox(MARGIN_L, SLIDE_H - Inches(1.0), CONTENT_W, Inches(0.4))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = tag
    _run(run, 13, bold=True, color=TEAL)
    _rect(slide, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), NAVY)
    nb = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.55), Inches(0.5), Inches(0.3))
    nb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    nr = nb.text_frame.paragraphs[0].add_run()
    nr.text = str(num)
    _run(nr, 11, color=WHITE)


def _slide(prs, label, title, bullets, tag, num, size=22):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), TEAL)

    lb = slide.shapes.add_textbox(MARGIN_L, Inches(0.6), CONTENT_W, Inches(0.35))
    lr = lb.text_frame.paragraphs[0].add_run()
    lr.text = label.upper()
    _run(lr, 11, bold=True, color=TEAL)

    tb = slide.shapes.add_textbox(MARGIN_L, Inches(1.15), CONTENT_W, Inches(0.9))
    tb.text_frame.word_wrap = True
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = title
    _run(tr, 32, bold=True, color=NAVY)

    bb = slide.shapes.add_textbox(MARGIN_L, Inches(2.2), CONTENT_W, Inches(3.8))
    tf = bb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        br = p.add_run()
        br.text = text
        _run(br, size, color=SLATE)

    _footer(slide, tag, num)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    n = 1

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    _rect(slide, Inches(0), Inches(5.6), SLIDE_W, Inches(0.08), TEAL)
    for y, text, sz, bold, col in [
        (1.8, "ROBOTICS · ROS 2", 14, True, TEAL),
        (2.5, "MoveIt2 Basics", 44, True, WHITE),
        (3.7, "Motion planning for robot arms & manipulators", 18, False, RGBColor(0xA0, 0xAE, 0xC0)),
        (6.0, "Educational Overview · 12 Slides", 13, True, ACCENT),
    ]:
        b = slide.shapes.add_textbox(MARGIN_L, Inches(y), CONTENT_W, Inches(0.8))
        run = b.text_frame.paragraphs[0].add_run()
        run.text = text
        _run(run, sz, bold=bold, color=col)
    n += 1

    slides = [
        (
            "Overview",
            "What Is MoveIt2?",
            [
                "Open-source motion planning framework for ROS 2.",
                "Plans safe arm trajectories — without hand-coding every move.",
                "Works with many industrial and research robots.",
            ],
            "MoveIt2 is the standard tool for robot arm motion in ROS 2.",
        ),
        (
            "Overview",
            "Why Use MoveIt2?",
            [
                "Handles collision avoidance automatically.",
                "Supports pick-and-place, welding, inspection, and more.",
                "Visualise and debug plans in RViz before running on hardware.",
            ],
            "It saves time and reduces risk when programming manipulators.",
        ),
        (
            "Concepts",
            "The Planning Scene",
            [
                "A virtual model of the robot's world.",
                "Includes the robot, table, obstacles, and objects to grasp.",
                "Updated as the environment changes.",
            ],
            "The planning scene is where MoveIt2 'sees' the workspace.",
        ),
        (
            "Concepts",
            "Robot Model & Planning Groups",
            [
                "URDF / SRDF files describe links, joints, and limits.",
                "A planning group is a set of joints — usually one arm.",
                "You plan motions for a group, not individual joints.",
            ],
            "Define your robot once; MoveIt2 handles the kinematics.",
        ),
        (
            "Concepts",
            "Motion Planning",
            [
                "Given a start pose and a goal pose, find a valid path.",
                "Planners (e.g. OMPL) search collision-free trajectories.",
                "Output: a sequence of joint positions over time.",
            ],
            "Planning = find a safe route; execution = follow that route.",
        ),
        (
            "Components",
            "Main Building Blocks",
            [
                "move_group — central node for planning & execution.",
                "Planning Scene Monitor — tracks robot + environment.",
                "MoveIt Setup Assistant — configure a new robot easily.",
            ],
            "Three pieces you will interact with most often.",
        ),
        (
            "Workflow",
            "Typical Use Flow",
            [
                "1. Define robot model (URDF / SRDF).",
                "2. Set up MoveIt2 with the Setup Assistant.",
                "3. Send a goal pose or joint target.",
                "4. Plan → preview in RViz → execute on the arm.",
            ],
            "Configure once, then plan and run motions from your application.",
        ),
        (
            "ROS 2",
            "How It Fits in ROS 2",
            [
                "Built on ROS 2 topics, services, and actions.",
                "Your node sends goals; MoveIt2 returns trajectories.",
                "Integrates with controllers via ros2_control.",
            ],
            "MoveIt2 is a layer above low-level motor control.",
        ),
        (
            "Applications",
            "Where MoveIt2 Is Used",
            [
                "Pick-and-place on factory lines.",
                "Bin picking and machine tending.",
                "Research labs — mobile manipulators, dual-arm setups.",
            ],
            "Any task that needs a robot arm to move intelligently.",
        ),
        (
            "Getting Started",
            "First Steps to Try",
            [
                "Install: ros-humble-moveit (or your ROS 2 distro).",
                "Run a demo: Panda arm in RViz.",
                "Use MoveIt Setup Assistant for your own robot.",
            ],
            "Start with the official tutorials — hands-on beats reading alone.",
        ),
        (
            "Wrap-Up",
            "Key Takeaways",
            [
                "MoveIt2 plans collision-free arm motions in ROS 2.",
                "Planning scene + robot model + planner = the core idea.",
                "Configure once, then focus on your application logic.",
            ],
            "MoveIt2 · Motion planning made accessible for education and industry.",
        ),
    ]

    for label, title, bullets, tag in slides:
        _slide(prs, label, title, bullets, tag, n)
        n += 1

    out = "/home/jjzzzz/Downloads/MoveIt2-Basics.pptx"
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    build()
